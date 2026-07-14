"""
agents/exporter.py — mechanical export of a node/artifact into a target
file format. No LLM call here, by design: same discipline as
agents/file_manager.py — deterministic and auditable, not generative.

Every exporter in this module takes the SAME input shape (the "common
artifact shape" from Part 0 §0.5), regardless of which domain produced
it:

    {
        "title": str,
        "sections": [
            {"heading": str, "content": str, "node_refs": [node_id, ...]},
            ...
        ],
        "metadata": {  # optional, all keys optional
            "workspace_id": str,
            "tags": [str, ...],
            "created_by": str,
            "created_at": str,
        },
    }

Notes/Research/Plan/etc. do not get their own exporter. A domain's report
generator, PRD writer, or note just needs one small adapter that shapes
its own object into the above before calling export_artifact() — see
graph/adapters.py for examples of that adapter layer.

Place this file at: agents/exporter.py
"""

import os
import csv
import json
import re
import datetime as _dt

from docx import Document
from docx.shared import Pt

from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt

from openpyxl import Workbook

from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.units import inch


SUPPORTED_FORMATS = ("docx", "pptx", "xlsx", "csv", "pdf", "md", "json")


# ---------------------------------------------------------------------------
# Shared validation / path safety (same pattern as file_manager.py's
# _safe_relpath — an exporter has no business writing outside the root
# it was told to write into).
# ---------------------------------------------------------------------------

def _validate_artifact(artifact: dict) -> dict:
    """Fills in defaults rather than raising on minor omissions — an
    artifact with a title and no sections (or vice versa) is still
    something worth exporting to *some* format, so this only raises on
    the one thing that would make every exporter below fail identically:
    a non-dict input.
    """
    if not isinstance(artifact, dict):
        raise ValueError("artifact must be a dict shaped like "
                          "{title, sections: [{heading, content, node_refs}]}")
    artifact = dict(artifact)
    artifact.setdefault("title", "Untitled")
    artifact.setdefault("sections", [])
    artifact.setdefault("metadata", {})
    normalized_sections = []
    for s in artifact["sections"]:
        normalized_sections.append({
            "heading": s.get("heading", ""),
            "content": s.get("content", "") or "",
            "node_refs": s.get("node_refs", []) or [],
        })
    artifact["sections"] = normalized_sections
    return artifact


def _safe_output_path(output_dir: str, filename: str) -> str:
    """Confines the write target to output_dir, same reasoning as
    file_manager.py's _safe_relpath — a filename derived from a
    user-editable title (see _slugify_filename) must never be allowed
    to escape the intended export directory via '../' tricks.
    """
    os.makedirs(output_dir, exist_ok=True)
    full = os.path.normpath(os.path.join(output_dir, filename))
    root = os.path.normpath(output_dir)
    if full != root and not full.startswith(root + os.sep):
        raise ValueError(f"Rejected unsafe export path: {filename}")
    return full


def _slugify_filename(title: str, ext: str, max_len: int = 60) -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "_", title.strip().lower()).strip("_")
    slug = slug[:max_len] or "untitled"
    return f"{slug}.{ext}"


# ---------------------------------------------------------------------------
# Format writers — each takes (artifact, output_path) and returns nothing;
# the dispatcher below handles path construction so these stay simple.
# ---------------------------------------------------------------------------

def _write_docx(artifact: dict, output_path: str) -> None:
    doc = Document()
    title_heading = doc.add_heading(artifact["title"], level=0)
    for section in artifact["sections"]:
        if section["heading"]:
            doc.add_heading(section["heading"], level=1)
        for para in (section["content"] or "").split("\n\n"):
            para = para.strip()
            if para:
                doc.add_paragraph(para)
        if section["node_refs"]:
            note = doc.add_paragraph()
            run = note.add_run("Sources: " + ", ".join(section["node_refs"]))
            run.italic = True
            run.font.size = Pt(9)
    doc.save(output_path)


def _write_pptx(artifact: dict, output_path: str) -> None:
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    body_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = artifact["title"]
    if artifact["sections"]:
        title_slide.placeholders[1].text = (
            f"{len(artifact['sections'])} section(s)"
        )

    for section in artifact["sections"]:
        slide = prs.slides.add_slide(body_layout)
        slide.shapes.title.text = section["heading"] or artifact["title"]
        body = slide.placeholders[1].text_frame
        body.clear()
        lines = [ln for ln in (section["content"] or "").split("\n") if ln.strip()]
        if not lines:
            lines = [""]
        body.text = lines[0]
        for line in lines[1:]:
            p = body.add_paragraph()
            p.text = line
        if section["node_refs"]:
            p = body.add_paragraph()
            p.text = "Sources: " + ", ".join(section["node_refs"])
            p.font.size = PptxPt(10)
            p.font.italic = True
    prs.save(output_path)


def _write_xlsx(artifact: dict, output_path: str) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = artifact["title"][:31] or "Sheet1"  # Excel's 31-char sheet-name limit
    ws.append(["heading", "content", "node_refs"])
    for section in artifact["sections"]:
        ws.append([
            section["heading"],
            section["content"],
            ", ".join(section["node_refs"]),
        ])
    wb.save(output_path)


def _write_csv(artifact: dict, output_path: str) -> None:
    with open(output_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["heading", "content", "node_refs"])
        for section in artifact["sections"]:
            writer.writerow([
                section["heading"],
                section["content"],
                ", ".join(section["node_refs"]),
            ])


def _write_pdf(artifact: dict, output_path: str) -> None:
    # Export-only, per Part 0 §0.5 — PDF import is a Notes-domain
    # ingestion concern (Part 4), not this module's job.
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(output_path, pagesize=LETTER,
                             topMargin=inch, bottomMargin=inch)
    flow = [Paragraph(artifact["title"], styles["Title"]), Spacer(1, 16)]
    for section in artifact["sections"]:
        if section["heading"]:
            flow.append(Paragraph(section["heading"], styles["Heading2"]))
            flow.append(Spacer(1, 6))
        for para in (section["content"] or "").split("\n\n"):
            para = para.strip()
            if para:
                # reportlab's Paragraph treats bare text as XML-ish markup;
                # escape the handful of characters that matter so content
                # containing "<", ">", "&" doesn't break rendering.
                escaped = (para.replace("&", "&amp;")
                               .replace("<", "&lt;")
                               .replace(">", "&gt;"))
                flow.append(Paragraph(escaped, styles["BodyText"]))
                flow.append(Spacer(1, 6))
        if section["node_refs"]:
            flow.append(Paragraph(
                "<i>Sources: " + ", ".join(section["node_refs"]) + "</i>",
                styles["BodyText"],
            ))
            flow.append(Spacer(1, 10))
    doc.build(flow)


def _write_md(artifact: dict, output_path: str) -> None:
    lines = [f"# {artifact['title']}", ""]
    for section in artifact["sections"]:
        if section["heading"]:
            lines.append(f"## {section['heading']}")
            lines.append("")
        if section["content"]:
            lines.append(section["content"])
            lines.append("")
        if section["node_refs"]:
            lines.append(f"*Sources: {', '.join(section['node_refs'])}*")
            lines.append("")
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines).rstrip() + "\n")


def _write_json(artifact: dict, output_path: str) -> None:
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(artifact, f, indent=2, ensure_ascii=False)


_WRITERS = {
    "docx": _write_docx,
    "pptx": _write_pptx,
    "xlsx": _write_xlsx,
    "csv": _write_csv,
    "pdf": _write_pdf,
    "md": _write_md,
    "json": _write_json,
}


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def export_artifact(artifact: dict, fmt: str, output_dir: str,
                     filename: str = None) -> str:
    """Mechanically writes `artifact` (the common {title, sections} shape)
    to `output_dir` in the requested format. Returns the full path written.

    This is the one function every domain (Notes, Research, Plan, ...)
    calls — no domain writes its own DOCX/PPTX/etc. writer. A domain's
    generator output just needs shaping into the common artifact form
    first (see graph/adapters.py) via a small per-domain adapter.
    """
    fmt = fmt.lower().lstrip(".")
    if fmt not in SUPPORTED_FORMATS:
        raise ValueError(f"Unsupported export format '{fmt}'. "
                          f"Supported: {', '.join(SUPPORTED_FORMATS)}")

    artifact = _validate_artifact(artifact)
    filename = filename or _slugify_filename(artifact["title"], fmt)
    if not filename.endswith(f".{fmt}"):
        filename = f"{filename}.{fmt}"

    output_path = _safe_output_path(output_dir, filename)
    _WRITERS[fmt](artifact, output_path)
    return output_path


if __name__ == "__main__":
    demo = {
        "title": "Demo Export",
        "sections": [
            {"heading": "Overview", "content": "First paragraph.\n\nSecond paragraph.",
             "node_refs": ["node:ws1:abc123"]},
            {"heading": "Details", "content": "Line one\nLine two", "node_refs": []},
        ],
        "metadata": {"workspace_id": "ws1"},
    }
    for f in SUPPORTED_FORMATS:
        path = export_artifact(demo, f, "/tmp/export_demo")
        print(f"wrote {path}")


# ---------------------------------------------------------------------------
# Part 6 §6.4 — Content calendar export. Deliberately NOT built on top of
# export_artifact()/the common {title, sections} artifact shape above:
# content_calendar_builder's real output is a structured
# {date, platform, content_ref} row list, not prose sections with
# node_refs, so squeezing it into the generic shape would either lose the
# structure or force a fake single "section" with the rows crammed into
# free text. Same deterministic, structured-in/file-out discipline as
# every writer above (no LLM call here either) — just a second, small,
# parallel entry point for a genuinely different input shape.
# ---------------------------------------------------------------------------

CALENDAR_FORMATS = ("ics", "csv", "md")

# Recognizes content_calendar_builder's own documented fallback labels
# (its seed brief in eo/registry.py: "day of launch", "day 3", "week 2",
# etc.) so the .ics writer below can still place a relative-sequencing
# calendar on a real timeline instead of silently refusing to export it.
# CSV/Markdown don't need this -- they show whatever date string was
# given, verbatim, since neither format requires a machine-parseable
# date.
_RELATIVE_DATE_RE = re.compile(
    r"^\s*(?:day\s+of\s+launch|day\s*0)\s*$|^\s*day\s+(\d+)\s*$|^\s*week\s+(\d+)\s*$",
    re.IGNORECASE,
)


def _validate_calendar_entries(entries: list) -> list:
    """Fills in defaults the same permissive way _validate_artifact()
    does above -- a row missing a platform or content_ref is still worth
    exporting, so this only raises on a non-list input, the one thing
    that would make every writer below fail identically."""
    if not isinstance(entries, list):
        raise ValueError("entries must be a list of "
                          "{date, platform, content_ref} dicts")
    normalized = []
    for e in entries:
        e = e or {}
        normalized.append({
            "date": str(e.get("date", "") or "").strip(),
            "platform": str(e.get("platform", "") or "").strip(),
            "content_ref": str(e.get("content_ref", "") or "").strip(),
        })
    return normalized


def _resolve_calendar_date(date_str: str, row_index: int,
                            anchor: _dt.date = None) -> _dt.date:
    """Best-effort mapping of one row's date field to a real calendar
    date, for the .ics writer only (CSV/Markdown show date_str verbatim
    and never call this).

    Tries, in order:
      1. A real ISO date ("2026-07-20") -- used exactly as given.
      2. content_calendar_builder's own relative labels ("day of
         launch"/"day 0", "day N", "week N") -- resolved as an offset
         from `anchor` (defaults to today, the date this file is being
         exported, since that's the only reasonable anchor a mechanical
         exporter has for a launch date it was never given).
      3. Anything else unparseable -- falls back to `anchor` + row_index
         days, purely so every row still lands on a distinct, ascending
         date in the .ics file rather than all stacking on one day.
    """
    anchor = anchor or _dt.date.today()
    try:
        return _dt.date.fromisoformat(date_str)
    except ValueError:
        pass
    m = _RELATIVE_DATE_RE.match(date_str)
    if m:
        if m.group(1):
            return anchor + _dt.timedelta(days=int(m.group(1)))
        if m.group(2):
            return anchor + _dt.timedelta(days=int(m.group(2)) * 7)
        return anchor  # "day of launch" / "day 0"
    return anchor + _dt.timedelta(days=row_index)


def _ics_escape(text: str) -> str:
    """RFC 5545 §3.3.11 text escaping -- backslash, semicolon, comma,
    then literal newlines, in that order (escaping the backslash first
    matters, or a later-inserted backslash would get re-escaped)."""
    return (text.replace("\\", "\\\\")
                .replace(";", "\\;")
                .replace(",", "\\,")
                .replace("\n", "\\n"))


def _write_calendar_ics(entries: list, output_path: str, title: str) -> None:
    """Hand-written ICS text via the stdlib only -- no icalendar PyPI
    dependency, matching this module's existing "deterministic and
    auditable" discipline. One all-day VEVENT per row; relative-label
    rows are placed via _resolve_calendar_date() above rather than
    omitted, so a calendar built without a real launch date (Part 6
    §6.4's documented fallback) still produces a genuinely importable
    .ics instead of an empty or invalid one."""
    anchor = _dt.date.today()
    lines = [
        "BEGIN:VCALENDAR",
        "VERSION:2.0",
        "PRODID:-//Growth Domain//Content Calendar//EN",
        "CALSCALE:GREGORIAN",
    ]
    stamp = _dt.datetime.now(_dt.timezone.utc).strftime("%Y%m%dT%H%M%SZ")
    for i, row in enumerate(entries):
        event_date = _resolve_calendar_date(row["date"], i, anchor)
        dtstart = event_date.strftime("%Y%m%d")
        dtend = (event_date + _dt.timedelta(days=1)).strftime("%Y%m%d")
        summary = _ics_escape(f"{row['platform'] or 'content'} — {title}")
        description_bits = []
        if row["content_ref"]:
            description_bits.append(f"Content: {row['content_ref']}")
        if row["date"] and row["date"] != event_date.isoformat():
            description_bits.append(f"Original schedule label: {row['date']}")
        lines += [
            "BEGIN:VEVENT",
            f"UID:{i}-{stamp}@growth-domain-calendar",
            f"DTSTAMP:{stamp}",
            f"DTSTART;VALUE=DATE:{dtstart}",
            f"DTEND;VALUE=DATE:{dtend}",
            f"SUMMARY:{summary}",
        ]
        if description_bits:
            lines.append(f"DESCRIPTION:{_ics_escape(chr(10).join(description_bits))}")
        lines.append("END:VEVENT")
    lines.append("END:VCALENDAR")
    # RFC 5545 §3.1 requires CRLF line endings.
    with open(output_path, "w", encoding="utf-8", newline="") as f:
        f.write("\r\n".join(lines) + "\r\n")


def _write_calendar_csv(entries: list, output_path: str, title: str) -> None:
    with open(output_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["date", "platform", "content_ref"])
        for row in entries:
            writer.writerow([row["date"], row["platform"], row["content_ref"]])


def _write_calendar_md(entries: list, output_path: str, title: str) -> None:
    lines = [f"# {title}", "", "| Date | Platform | Content |",
             "|------|----------|---------|"]
    for row in entries:
        # Pipe characters in a cell would break the table -- escape same
        # as _write_md() would if it ever needed table cells (it doesn't,
        # today), so a content_ref containing "|" can't corrupt the row.
        cells = [c.replace("|", "\\|") for c in
                 (row["date"], row["platform"], row["content_ref"])]
        lines.append(f"| {cells[0]} | {cells[1]} | {cells[2]} |")
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines).rstrip() + "\n")


_CALENDAR_WRITERS = {
    "ics": _write_calendar_ics,
    "csv": _write_calendar_csv,
    "md": _write_calendar_md,
}


def export_content_calendar(entries: list, fmt: str, output_dir: str,
                             title: str = "Content Calendar",
                             filename: str = None) -> str:
    """The content-calendar counterpart to export_artifact() above.

    entries: the structured {date, platform, content_ref} list
        content_calendar_builder produces (Part 6 §6.4) -- date may be a
        real ISO date OR one of its documented relative-sequencing
        fallback labels ("day of launch", "day 3", "week 2") when no
        real launch date was available upstream.
    fmt: one of CALENDAR_FORMATS ("ics", "csv", "md").
    Returns the full path written.
    """
    fmt = fmt.lower().lstrip(".")
    if fmt not in CALENDAR_FORMATS:
        raise ValueError(f"Unsupported calendar export format '{fmt}'. "
                          f"Supported: {', '.join(CALENDAR_FORMATS)}")

    entries = _validate_calendar_entries(entries)
    filename = filename or _slugify_filename(title, fmt)
    if not filename.endswith(f".{fmt}"):
        filename = f"{filename}.{fmt}"

    output_path = _safe_output_path(output_dir, filename)
    _CALENDAR_WRITERS[fmt](entries, output_path, title)
    return output_path


if __name__ == "__main__":
    demo_calendar = [
        {"date": "day of launch", "platform": "twitter", "content_ref": "Launch announcement (short)"},
        {"date": "day 3", "platform": "linkedin", "content_ref": "Launch announcement (long-form)"},
        {"date": "week 2", "platform": "press_release", "content_ref": "Press release"},
    ]
    for f in CALENDAR_FORMATS:
        path = export_content_calendar(demo_calendar, f, "/tmp/export_demo")
        print(f"wrote {path}")