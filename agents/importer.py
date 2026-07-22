"""
agents/importer.py — mechanical import of docx/pptx/xlsx/csv/md/json back
into the common artifact shape ({title, sections: [{heading, content,
node_refs}]}) defined in agents/exporter.py. No LLM call here either —
same deterministic, auditable discipline as the exporter and as
agents/file_manager.py.

PDF import is deliberately NOT here — per Part 0 §0.5, PDF import is a
Notes-domain ingestion concern (Part 4: OCR, layout-aware extraction,
etc. — a fundamentally different, heavier problem than "parse structured
text back out of a format we ourselves wrote"). Everything below assumes
reasonably well-structured input, which holds for anything this system's
own exporter produced, and for hand-authored docx/pptx/xlsx/md/json that
follows normal heading/section conventions.

Place this file at: agents/importer.py
"""

import os
import csv
import json

from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


SUPPORTED_FORMATS = ("docx", "pptx", "xlsx", "csv", "md", "json")


def _new_section(heading: str = "") -> dict:
    return {"heading": heading, "content": "", "node_refs": []}


def _append_content(section: dict, text: str) -> None:
    text = text.strip()
    if not text:
        return
    section["content"] = (section["content"] + "\n\n" + text).strip() \
        if section["content"] else text


# ---------------------------------------------------------------------------
# Format readers — each takes a file path and returns the common shape.
# ---------------------------------------------------------------------------

def _read_docx(path: str, default_title: str = None) -> dict:
    doc = Document(path)
    title = default_title or os.path.splitext(os.path.basename(path))[0]
    sections = []
    current = None

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower()

        if style == "title" or style == "heading 0":
            title = text
        elif style.startswith("heading"):
            # Any heading level below Title starts a new section — mirrors
            # the exporter, which only ever emits level-1 section headings,
            # but hand-authored docs may use level 1 or 2 for structure.
            current = _new_section(text)
            sections.append(current)
        else:
            if text.startswith("Sources: ") and current is not None:
                refs = text[len("Sources: "):].split(", ")
                current["node_refs"] = [r for r in refs if r]
                continue
            if current is None:
                current = _new_section()
                sections.append(current)
            _append_content(current, text)

    return {"title": title, "sections": sections, "metadata": {}}


def _read_pptx(path: str, default_title: str = None) -> dict:
    prs = Presentation(path)
    slides = list(prs.slides)
    title = default_title or os.path.splitext(os.path.basename(path))[0]
    sections = []

    for i, slide in enumerate(slides):
        texts = []
        slide_title = None
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape == slide.shapes.title:
                slide_title = text
            else:
                texts.append(text)

        if i == 0:
            # First slide is always the title slide, matching exactly how
            # the exporter builds it (title_layout + optional "N
            # section(s)" subtitle) — never a content section.
            if slide_title:
                title = slide_title
            continue

        section = _new_section(slide_title or "")
        node_refs = []
        body_lines = []
        for text in texts:
            for line in text.split("\n"):
                line = line.strip()
                if not line:
                    continue
                if line.startswith("Sources: "):
                    node_refs.extend(r for r in line[len("Sources: "):].split(", ") if r)
                else:
                    body_lines.append(line)
        section["content"] = "\n".join(body_lines)
        section["node_refs"] = node_refs
        sections.append(section)

    return {"title": title, "sections": sections, "metadata": {}}


def _read_xlsx(path: str) -> dict:
    wb = load_workbook(path, read_only=True)
    ws = wb.active
    title = ws.title
    sections = []
    rows = list(ws.iter_rows(values_only=True))
    if rows and rows[0] and str(rows[0][0]).lower() == "heading":
        rows = rows[1:]  # skip the header row the exporter writes
    for row in rows:
        if not row or all(c is None for c in row):
            continue
        heading = row[0] if len(row) > 0 and row[0] else ""
        content = row[1] if len(row) > 1 and row[1] else ""
        refs_raw = row[2] if len(row) > 2 and row[2] else ""
        node_refs = [r.strip() for r in str(refs_raw).split(",") if r.strip()]
        sections.append({"heading": heading, "content": content, "node_refs": node_refs})
    return {"title": title, "sections": sections, "metadata": {}}


def _read_csv(path: str, default_title: str = None) -> dict:
    title = default_title or os.path.splitext(os.path.basename(path))[0]
    sections = []
    with open(path, newline="", encoding="utf-8") as f:
        reader = csv.reader(f)
        rows = list(reader)
    if rows and rows[0] and rows[0][0].lower() == "heading":
        rows = rows[1:]
    for row in rows:
        if not row or not any(row):
            continue
        heading = row[0] if len(row) > 0 else ""
        content = row[1] if len(row) > 1 else ""
        refs_raw = row[2] if len(row) > 2 else ""
        node_refs = [r.strip() for r in refs_raw.split(",") if r.strip()]
        sections.append({"heading": heading, "content": content, "node_refs": node_refs})
    return {"title": title, "sections": sections, "metadata": {}}


def parse_markdown_text(text: str, default_title: str = "Untitled") -> dict:
    """The actual markdown -> {title, sections, metadata} parser. Split
    out of _read_md() below so Part 4 §4.4's generator roles (mapper,
    report_writer, slide_planner, podcast_scriptwriter — every one of
    which asks generic_worker for headered Markdown via its
    MARKDOWN_INSTRUCTION) can turn their raw output straight into an
    exportable artifact without going through a temp file first. This
    is the only place that understands the heading/"Sources:" grammar;
    _read_md() and graph/adapters.py's markdown_text_to_artifact() both
    call this rather than re-implementing it.

    Parses blank-line-delimited blocks rather than individual lines.
    This matters because _write_md() writes a section's whole `content`
    string as a single block (which may itself contain internal '\\n\\n'
    paragraph breaks) — parsing line-by-line would insert an extra blank
    line between every line of multi-line content. Splitting on blank
    lines and re-joining multiple content blocks with '\\n\\n' exactly
    inverts that, including the case where a content block's own internal
    '\\n\\n' causes it to split into more than one block here.
    """
    blocks = [b for b in text.split("\n\n")]
    title = default_title
    sections = []
    current = None

    for block in blocks:
        stripped = block.strip()
        if not stripped:
            continue
        if stripped.startswith("# "):
            title = stripped[2:].strip()
        elif stripped.startswith("## "):
            current = _new_section(stripped[3:].strip())
            sections.append(current)
        elif stripped.startswith("*Sources: ") and stripped.endswith("*") and current is not None:
            refs = stripped[len("*Sources: "):-1].split(", ")
            current["node_refs"] = [r for r in refs if r]
        else:
            if current is None:
                current = _new_section()
                sections.append(current)
            _append_content(current, block.strip("\n"))

    return {"title": title, "sections": sections, "metadata": {}}


def _read_md(path: str, default_title: str = None) -> dict:
    with open(path, encoding="utf-8") as f:
        text = f.read()
    default_title = default_title or os.path.splitext(os.path.basename(path))[0]
    return parse_markdown_text(text, default_title)


def _read_json(path: str) -> dict:
    with open(path, encoding="utf-8") as f:
        data = json.load(f)
    if not isinstance(data, dict) or "sections" not in data:
        raise ValueError(f"{path} is not a valid artifact JSON "
                          "(expected {title, sections: [...]})")
    data.setdefault("title", "Untitled")
    data.setdefault("metadata", {})
    return data


_READERS = {
    "docx": _read_docx,
    "pptx": _read_pptx,
    "xlsx": _read_xlsx,
    "csv": _read_csv,
    "md": _read_md,
    "json": _read_json,
}


def import_artifact(path: str, fmt: str = None, default_title: str = None) -> dict:
    """Mechanically parses a file at `path` back into the common
    {title, sections, metadata} shape. `fmt` is inferred from the file
    extension when not given explicitly.

    `default_title`: the title fallback to use when the file itself has
    no in-file title (e.g. no "Title"-styled paragraph in a docx, no "#
    Heading" in markdown). Callers that read from a temp upload path
    (api/server.py's /api/notes/import endpoint) should pass the
    original uploaded filename here -- otherwise the fallback silently
    becomes the temp file's random name (e.g. "tmp7nfj4s2h") instead of
    something the user recognizes. When omitted, falls back to `path`'s
    own basename, same as before.
    """
    if fmt is None:
        fmt = os.path.splitext(path)[1].lstrip(".").lower()
    else:
        fmt = fmt.lower().lstrip(".")

    if fmt not in SUPPORTED_FORMATS:
        raise ValueError(f"Unsupported import format '{fmt}'. "
                          f"Supported: {', '.join(SUPPORTED_FORMATS)}. "
                          "(PDF import is a Notes-domain concern, not this module.)")
    if not os.path.exists(path):
        raise FileNotFoundError(path)

    reader = _READERS[fmt]
    if fmt in ("xlsx", "json"):
        # Neither reader's title comes from the temp path (xlsx uses the
        # sheet name, json uses/rejects its own "title" field), so
        # there's no random-name fallback to override here.
        return reader(path)
    return reader(path, default_title=default_title)


if __name__ == "__main__":
    import sys
    for p in sys.argv[1:]:
        artifact = import_artifact(p)
        print(f"--- {p} ---")
        print(json.dumps(artifact, indent=2)[:500])